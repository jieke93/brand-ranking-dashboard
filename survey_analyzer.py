"""
설문조사 RAW 데이터 분석 도구
- 엑셀 RAW 파일을 읽어서 나이/성별 기준으로 평균값을 산출
- 아이템별로 선호도, 적정가격 등의 평균을 계산
- 결과를 깔끔한 엑셀 파일로 출력
"""

import sys
import os
import re
from collections import defaultdict
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


def clean_gender(raw_gender):
    """성별 라벨을 '남성'/'여성'으로 정리"""
    g = str(raw_gender).strip()
    if g.startswith("남성"):
        return "남성"
    elif g.startswith("여성"):
        return "여성"
    return g


def load_raw_data(filepath):
    """RAW 엑셀 파일을 로드하고 헤더와 데이터를 반환"""
    wb = openpyxl.load_workbook(filepath, data_only=True)
    ws = wb.active

    headers = []
    for col in range(1, ws.max_column + 1):
        val = ws.cell(row=1, column=col).value
        headers.append(val if val else "")

    data = []
    for row in range(2, ws.max_row + 1):
        row_data = []
        for col in range(1, ws.max_column + 1):
            row_data.append(ws.cell(row=row, column=col).value)
        # 빈 행 스킵
        if any(v is not None for v in row_data[:6]):
            data.append(row_data)

    return headers, data


def find_gender_age_columns(headers):
    """성별과 연령 열 인덱스를 찾기"""
    gender_col = None
    age_col = None
    for i, h in enumerate(headers):
        h_str = str(h).strip()
        if "성별" in h_str:
            gender_col = i
        if "연령" in h_str or "나이" in h_str:
            age_col = i
    return gender_col, age_col


def identify_items(headers, start_col=6):
    """
    G열(인덱스 6)부터 시작하여 아이템 경계를 찾음
    '1.'로 시작하는 헤더가 나오면 새 아이템의 시작
    각 아이템은 {name, questions: [{col_idx, header, type}]} 형태
    """
    items = []
    current_item = None
    item_counter = 0

    for i in range(start_col, len(headers)):
        h = str(headers[i]).strip()
        if not h:
            continue
        # 빈 Google Sheets 열 ("열1", "열2" 등) 무시
        if re.match(r'^열\d+$', h):
            continue

        # "1. " 또는 "1.[" 로 시작하면 새 아이템
        if re.match(r'^1\.[\s\[]', h):
            if current_item:
                items.append(current_item)
            item_counter += 1
            current_item = {
                "item_no": item_counter,
                "questions": []
            }
            current_item["questions"].append({
                "col_idx": i,
                "header": h,
                "type": classify_question(h)
            })
        elif current_item is not None:
            q_type = classify_question(h)
            # 숫자로 시작하는 질문(2., 3., 4. 등)이거나 특수질문
            current_item["questions"].append({
                "col_idx": i,
                "header": h,
                "type": q_type
            })
        else:
            # 아이템 시작 전 특수 질문 → 무시하거나 별도 처리
            pass

    if current_item:
        items.append(current_item)

    return items


def classify_question(header):
    """질문 유형 분류"""
    h = str(header).strip()
    if "선호도" in h and ("10점" in h or "만점" in h):
        return "선호도"
    elif "가격" in h or "적정" in h:
        return "적정가격"
    elif "컬러" in h and ("구매의향" in h or "골라" in h or "구매 의향" in h):
        return "컬러선호"
    elif "기장" in h:
        return "기장비교"
    elif "구매하고 싶은" in h or "구매하고싶은" in h:
        return "구매희망"
    elif "자유롭게" in h or "만족/불만족" in h or "제안" in h:
        return "주관식"
    elif "선호하는" in h and ("유형" in h or "디자인" in h):
        return "선호유형"
    elif "컬러 선호도" in h:
        return "컬러선호도"
    elif re.match(r'^\d+\.', h):
        return "기타수치"
    else:
        return "기타"


def safe_numeric(val):
    """값을 숫자로 변환 시도"""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        return float(val)
    val_str = str(val).strip().replace(",", "").replace("원", "").replace("만", "0000")
    # 숫자만 추출
    val_str = re.sub(r'[^\d.]', '', val_str)
    if val_str:
        try:
            return float(val_str)
        except ValueError:
            return None
    return None


def calculate_averages(data, col_idx, gender_col, age_col, exclude_zero=False):
    """특정 열의 평균값을 전체/성별/나이그룹별로 계산
    exclude_zero=True이면 0값을 제외하고 계산 (구매의사 없음 제외)"""
    results = {
        "전체": [],
        "성별": defaultdict(list),
        "연령": defaultdict(list),
        "성별_연령": defaultdict(list),
    }

    for row in data:
        val = safe_numeric(row[col_idx])
        if val is None:
            continue
        if exclude_zero and val == 0:
            continue

        gender = clean_gender(row[gender_col]) if gender_col is not None and row[gender_col] else "미응답"
        age = str(row[age_col]).strip() if age_col is not None and row[age_col] else "미응답"

        results["전체"].append(val)
        results["성별"][gender].append(val)
        results["연령"][age].append(val)
        results["성별_연령"][f"{gender}_{age}"].append(val)

    # 평균 계산
    avg = {}
    avg["전체"] = round(sum(results["전체"]) / len(results["전체"]), 1) if results["전체"] else "-"
    avg["전체_응답수"] = len(results["전체"])

    avg["성별"] = {}
    for k, v in results["성별"].items():
        avg["성별"][k] = {"평균": round(sum(v) / len(v), 1) if v else "-", "응답수": len(v)}

    avg["연령"] = {}
    for k, v in results["연령"].items():
        avg["연령"][k] = {"평균": round(sum(v) / len(v), 1) if v else "-", "응답수": len(v)}

    avg["성별_연령"] = {}
    for k, v in results["성별_연령"].items():
        avg["성별_연령"][k] = {"평균": round(sum(v) / len(v), 1) if v else "-", "응답수": len(v)}

    return avg


def collect_color_stats(data, col_idx, gender_col, age_col):
    """컬러 선호 통계 수집"""
    results = {
        "전체": defaultdict(int),
        "성별": defaultdict(lambda: defaultdict(int)),
        "연령": defaultdict(lambda: defaultdict(int)),
    }
    total_respondents = 0

    for row in data:
        val = row[col_idx]
        if val is None or str(val).strip() in ("", "없음", "0"):
            continue

        gender = clean_gender(row[gender_col]) if gender_col is not None and row[gender_col] else "미응답"
        age = str(row[age_col]).strip() if age_col is not None and row[age_col] else "미응답"

        colors = [c.strip() for c in str(val).split(",") if c.strip()]
        total_respondents += 1
        for color in colors:
            results["전체"][color] += 1
            results["성별"][gender][color] += 1
            results["연령"][age][color] += 1

    return results, total_respondents


# 특수 질문 유형 (공통 질문이 아닌 것들)
COMMON_TYPES = {"선호도", "적정가격", "컬러선호"}


def collect_special_stats(data, col_idx, header, q_type):
    """특수 질문의 통계를 수집.
    - 수치형(컬러선호도, 기타수치): 평균/분포 반환
    - 선택형(기장비교, 선호유형, 구매희망): 빈도 분포 반환
    - 주관식: 빈도 분포 (상위 10개)
    """
    vals = []
    for row in data:
        v = row[col_idx]
        if v is None or str(v).strip() in ("", "없음"):
            continue
        vals.append(str(v).strip())

    if not vals:
        return {"type": "empty", "header": header, "q_type": q_type, "data": []}

    # 수치형 판단
    if q_type in ("컬러선호도", "기타수치"):
        nums = []
        for v in vals:
            try:
                nums.append(float(re.sub(r'[^\d.]', '', v)))
            except ValueError:
                pass
        if nums:
            avg_all = round(sum(nums) / len(nums), 1)
            nums_ex = [n for n in nums if n != 0]
            avg_ex = round(sum(nums_ex) / len(nums_ex), 1) if nums_ex else "-"
            return {
                "type": "numeric",
                "header": header,
                "q_type": q_type,
                "avg": avg_all,
                "avg_ex": avg_ex,
                "count": len(nums),
                "count_ex": len(nums_ex),
            }

    # 선택형/텍스트: 빈도 분포
    counter = defaultdict(int)
    for v in vals:
        # 쉼표로 복수 응답 분리
        parts = [p.strip() for p in v.split(",") if p.strip()]
        for p in parts:
            counter[p] += 1
    sorted_items = sorted(counter.items(), key=lambda x: -x[1])
    return {
        "type": "choice",
        "header": header,
        "q_type": q_type,
        "total": len(vals),
        "distribution": sorted_items,  # [(value, count), ...]
    }


def sort_age_groups(age_groups):
    """연령 그룹을 숫자 기준으로 정렬"""
    def age_sort_key(age_str):
        match = re.search(r'(\d+)', str(age_str))
        return int(match.group(1)) if match else 999
    return sorted(age_groups, key=age_sort_key)


def create_summary_excel(headers, data, items, gender_col, age_col, output_path):
    """분석 결과를 깔끔한 엑셀로 출력"""
    wb = openpyxl.Workbook()

    # 스타일 정의
    header_font = Font(bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    item_fill = PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid")
    item_font = Font(bold=True, size=11, color="1F3864")
    sub_header_fill = PatternFill(start_color="E9EDF4", end_color="E9EDF4", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin', color='B4C6E7'),
        right=Side(style='thin', color='B4C6E7'),
        top=Side(style='thin', color='B4C6E7'),
        bottom=Side(style='thin', color='B4C6E7')
    )
    center_align = Alignment(horizontal='center', vertical='center', wrap_text=True)
    left_align = Alignment(horizontal='left', vertical='center', wrap_text=True)

    # 성별/연령 목록 수집
    genders = set()
    ages = set()
    for row in data:
        g = clean_gender(row[gender_col]) if gender_col is not None and row[gender_col] else "미응답"
        a = str(row[age_col]).strip() if age_col is not None and row[age_col] else "미응답"
        genders.add(g)
        ages.add(a)
    genders = sorted(genders)
    ages = sort_age_groups(ages)

    # ==============================
    # 시트1: 선호도 요약
    # ==============================
    ws_pref = wb.active
    ws_pref.title = "선호도 요약"
    _write_numeric_summary(ws_pref, items, data, gender_col, age_col, genders, ages,
                           "선호도", "상품별 선호도 평균 (10점 만점)",
                           header_font, header_fill, item_fill, item_font, sub_header_fill,
                           thin_border, center_align, left_align)

    # ==============================
    # 시트2: 적정가격 요약
    # ==============================
    ws_price = wb.create_sheet("적정가격 요약")
    _write_numeric_summary(ws_price, items, data, gender_col, age_col, genders, ages,
                           "적정가격", "상품별 적정가격 평균 (원)",
                           header_font, header_fill, item_fill, item_font, sub_header_fill,
                           thin_border, center_align, left_align)

    # ==============================
    # 시트3: 컬러 선호 요약
    # ==============================
    ws_color = wb.create_sheet("컬러선호 요약")
    _write_color_summary(ws_color, items, data, gender_col, age_col, genders, ages,
                         header_font, header_fill, item_fill, item_font, sub_header_fill,
                         thin_border, center_align, left_align)

    # ==============================
    # 시트4: 성별x연령 교차분석
    # ==============================
    ws_cross = wb.create_sheet("성별x연령 교차분석")
    _write_cross_summary(ws_cross, items, data, gender_col, age_col, genders, ages,
                         header_font, header_fill, item_fill, item_font, sub_header_fill,
                         thin_border, center_align, left_align)

    # ==============================
    # 시트5: 데이터 개요
    # ==============================
    ws_overview = wb.create_sheet("데이터 개요")
    _write_overview(ws_overview, data, gender_col, age_col, genders, ages, items,
                    header_font, header_fill, thin_border, center_align, left_align)

    # ==============================
    # 시트6: 특수 질문 요약
    # ==============================
    ws_special = wb.create_sheet("특수질문 요약")
    _write_special_summary(ws_special, items, data,
                           header_font, header_fill, item_fill, item_font,
                           thin_border, center_align, left_align)

    # ==============================
    # 시트7: 아이템 대시보드 (숨겨진 데이터 시트 + 대시보드)
    # ==============================
    ws_db = wb.create_sheet("_DB")
    _write_dashboard_data(ws_db, items, data, gender_col, age_col, genders, ages)
    ws_db.sheet_state = 'hidden'

    ws_dash = wb.create_sheet("아이템 대시보드")
    _write_dashboard(ws_dash, items, data, gender_col, age_col, genders, ages,
                     header_font, header_fill, item_fill, item_font, sub_header_fill,
                     thin_border, center_align, left_align)
    # 대시보드를 맨 앞으로 이동
    wb.move_sheet(ws_dash, offset=-(len(wb.sheetnames) - 1))

    wb.save(output_path)
    print(f"\n✅ 분석 완료! 결과 파일: {output_path}")
    print(f"   - 총 응답자: {len(data)}명")
    print(f"   - 분석 아이템: {len(items)}개")
    print(f"   - 시트: 아이템 대시보드, 선호도 요약, 적정가격 요약, 컬러선호 요약, 성별x연령 교차분석, 특수질문 요약, 데이터 개요")


def _write_numeric_summary(ws, items, data, gender_col, age_col, genders, ages,
                           q_type, title, header_font, header_fill, item_fill, item_font,
                           sub_header_fill, thin_border, center_align, left_align):
    """선호도/적정가격 숫자 요약 시트 작성 (0포함/0제외 이중 표시)"""
    row = 1
    zero_excl_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")  # 노란 배경

    # 제목
    ws.cell(row=row, column=1, value=title).font = Font(bold=True, size=14, color="1F3864")
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=5 + len(genders) + len(ages))
    row += 2

    # 헤더 행
    col_headers = ["No", "구분", "아이템 질문", "전체평균", "응답수"]
    for g in genders:
        col_headers.append(f"{g}")
    for a in ages:
        col_headers.append(f"{a}")

    for c, h in enumerate(col_headers, 1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    row += 1

    # 데이터 행
    for item in items:
        numeric_qs = [q for q in item["questions"] if q["type"] == q_type]
        if not numeric_qs:
            continue

        for q in numeric_qs:
            # --- 0 포함 행 ---
            avg_all = calculate_averages(data, q["col_idx"], gender_col, age_col, exclude_zero=False)
            _write_avg_row(ws, row, item["item_no"], "0포함", q["header"][:70], avg_all,
                           genders, ages, thin_border, center_align, left_align, None)
            row += 1

            # --- 0 제외 행 ---
            avg_ex = calculate_averages(data, q["col_idx"], gender_col, age_col, exclude_zero=True)
            _write_avg_row(ws, row, "", "0제외", "", avg_ex,
                           genders, ages, thin_border, center_align, left_align, zero_excl_fill)
            row += 1

    # 열 너비 조정
    ws.column_dimensions['A'].width = 6
    ws.column_dimensions['B'].width = 8
    ws.column_dimensions['C'].width = 45
    ws.column_dimensions['D'].width = 12
    ws.column_dimensions['E'].width = 8
    for c in range(6, 6 + len(genders) + len(ages)):
        ws.column_dimensions[get_column_letter(c)].width = 12

    # 필터 적용
    ws.auto_filter.ref = f"A3:{get_column_letter(5 + len(genders) + len(ages))}{row - 1}"


def _write_avg_row(ws, row, item_no, label, q_text, avg, genders, ages,
                   thin_border, center_align, left_align, bg_fill):
    """평균 데이터 1행 기록 (0포함/0제외 공용)"""
    ws.cell(row=row, column=1, value=item_no).alignment = center_align
    ws.cell(row=row, column=1).border = thin_border

    ws.cell(row=row, column=2, value=label).alignment = center_align
    ws.cell(row=row, column=2).border = thin_border
    ws.cell(row=row, column=2).font = Font(bold=True, size=9,
                                            color="C00000" if label == "0제외" else "333333")

    ws.cell(row=row, column=3, value=q_text).alignment = left_align
    ws.cell(row=row, column=3).border = thin_border

    ws.cell(row=row, column=4, value=avg["전체"]).alignment = center_align
    ws.cell(row=row, column=4).border = thin_border
    ws.cell(row=row, column=4).font = Font(bold=True)

    ws.cell(row=row, column=5, value=avg["전체_응답수"]).alignment = center_align
    ws.cell(row=row, column=5).border = thin_border

    col = 6
    for g in genders:
        val = avg["성별"].get(g, {}).get("평균", "-")
        ws.cell(row=row, column=col, value=val).alignment = center_align
        ws.cell(row=row, column=col).border = thin_border
        col += 1

    for a in ages:
        val = avg["연령"].get(a, {}).get("평균", "-")
        ws.cell(row=row, column=col, value=val).alignment = center_align
        ws.cell(row=row, column=col).border = thin_border
        col += 1

    if bg_fill:
        for c in range(1, col):
            ws.cell(row=row, column=c).fill = bg_fill


def _write_color_summary(ws, items, data, gender_col, age_col, genders, ages,
                         header_font, header_fill, item_fill, item_font, sub_header_fill,
                         thin_border, center_align, left_align):
    """컬러 선호 요약 시트"""
    row = 1
    ws.cell(row=row, column=1, value="상품별 컬러 선호 요약").font = Font(bold=True, size=14, color="1F3864")
    row += 2

    for item in items:
        color_qs = [q for q in item["questions"] if q["type"] in ("컬러선호", "컬러선호도")]
        if not color_qs:
            continue

        for q in color_qs:
            # 아이템 제목
            ws.cell(row=row, column=1, value=f"아이템 {item['item_no']}").font = item_font
            ws.cell(row=row, column=1).fill = item_fill
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=5)
            for c in range(1, 6):
                ws.cell(row=row, column=c).fill = item_fill
                ws.cell(row=row, column=c).border = thin_border
            row += 1

            stats, total = collect_color_stats(data, q["col_idx"], gender_col, age_col)

            # 헤더
            for c, h in enumerate(["컬러", "응답수", "비율(%)", "성별분포", "연령분포"], 1):
                cell = ws.cell(row=row, column=c, value=h)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = center_align
                cell.border = thin_border
            row += 1

            # 컬러별 데이터 (응답수 내림차순)
            sorted_colors = sorted(stats["전체"].items(), key=lambda x: -x[1])
            for color, count in sorted_colors:
                ws.cell(row=row, column=1, value=color).alignment = left_align
                ws.cell(row=row, column=1).border = thin_border
                ws.cell(row=row, column=2, value=count).alignment = center_align
                ws.cell(row=row, column=2).border = thin_border
                pct = round(count / total * 100, 1) if total > 0 else 0
                ws.cell(row=row, column=3, value=pct).alignment = center_align
                ws.cell(row=row, column=3).border = thin_border

                # 성별 분포
                gender_parts = []
                for g in genders:
                    gc = stats["성별"][g].get(color, 0)
                    if gc > 0:
                        gender_parts.append(f"{g}:{gc}")
                ws.cell(row=row, column=4, value=", ".join(gender_parts)).alignment = left_align
                ws.cell(row=row, column=4).border = thin_border

                # 연령 분포
                age_parts = []
                for a in sort_age_groups(stats["연령"].keys()):
                    ac = stats["연령"][a].get(color, 0)
                    if ac > 0:
                        age_parts.append(f"{a}:{ac}")
                ws.cell(row=row, column=5, value=", ".join(age_parts)).alignment = left_align
                ws.cell(row=row, column=5).border = thin_border
                row += 1

            row += 1  # 아이템 간 간격

    ws.column_dimensions['A'].width = 18
    ws.column_dimensions['B'].width = 10
    ws.column_dimensions['C'].width = 10
    ws.column_dimensions['D'].width = 30
    ws.column_dimensions['E'].width = 50


def _write_cross_summary(ws, items, data, gender_col, age_col, genders, ages,
                         header_font, header_fill, item_fill, item_font, sub_header_fill,
                         thin_border, center_align, left_align):
    """성별x연령 교차분석 시트 (선호도 기준)"""
    row = 1
    ws.cell(row=row, column=1, value="성별 × 연령 교차분석 (선호도 평균)").font = Font(bold=True, size=14, color="1F3864")
    row += 2

    # 헤더
    col_headers = ["No", "아이템"]
    for g in genders:
        for a in ages:
            col_headers.append(f"{g}\n{a}")
    for c, h in enumerate(col_headers, 1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    row += 1

    for item in items:
        pref_qs = [q for q in item["questions"] if q["type"] == "선호도"]
        if not pref_qs:
            continue

        q = pref_qs[0]
        # 0포함 행
        avg = calculate_averages(data, q["col_idx"], gender_col, age_col, exclude_zero=False)
        ws.cell(row=row, column=1, value=item["item_no"]).alignment = center_align
        ws.cell(row=row, column=1).border = thin_border
        ws.cell(row=row, column=2, value=f"아이템{item['item_no']} (0포함)").alignment = left_align
        ws.cell(row=row, column=2).border = thin_border
        col = 3
        for g in genders:
            for a in ages:
                key = f"{g}_{a}"
                val = avg["성별_연령"].get(key, {}).get("평균", "-")
                cell = ws.cell(row=row, column=col, value=val)
                cell.alignment = center_align
                cell.border = thin_border
                if isinstance(val, (int, float)):
                    if val >= 8:
                        cell.fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
                    elif val >= 6:
                        cell.fill = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
                    elif val < 4:
                        cell.fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
                col += 1
        row += 1

        # 0제외 행
        avg_ex = calculate_averages(data, q["col_idx"], gender_col, age_col, exclude_zero=True)
        ws.cell(row=row, column=1, value="").border = thin_border
        ws.cell(row=row, column=2, value=f"아이템{item['item_no']} (0제외)").alignment = left_align
        ws.cell(row=row, column=2).border = thin_border
        ws.cell(row=row, column=2).font = Font(bold=True, size=9, color="C00000")
        col = 3
        zero_bg = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
        for g in genders:
            for a in ages:
                key = f"{g}_{a}"
                val = avg_ex["성별_연령"].get(key, {}).get("평균", "-")
                cell = ws.cell(row=row, column=col, value=val)
                cell.alignment = center_align
                cell.border = thin_border
                cell.fill = zero_bg
                col += 1
        for c in range(1, 3):
            ws.cell(row=row, column=c).fill = zero_bg
        row += 1

    ws.column_dimensions['A'].width = 6
    ws.column_dimensions['B'].width = 15
    for c in range(3, 3 + len(genders) * len(ages)):
        ws.column_dimensions[get_column_letter(c)].width = 12


def _write_overview(ws, data, gender_col, age_col, genders, ages, items,
                    header_font, header_fill, thin_border, center_align, left_align):
    """데이터 개요 시트"""
    row = 1
    ws.cell(row=row, column=1, value="설문 데이터 개요").font = Font(bold=True, size=14, color="1F3864")
    row += 2

    ws.cell(row=row, column=1, value="총 응답자 수").font = Font(bold=True)
    ws.cell(row=row, column=2, value=len(data))
    row += 1
    ws.cell(row=row, column=1, value="분석 아이템 수").font = Font(bold=True)
    ws.cell(row=row, column=2, value=len(items))
    row += 2

    # 성별 분포
    ws.cell(row=row, column=1, value="성별 분포").font = Font(bold=True, size=12)
    row += 1
    for c, h in enumerate(["성별", "인원", "비율(%)"], 1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    row += 1

    gender_counts = defaultdict(int)
    for r in data:
        g = clean_gender(r[gender_col]) if gender_col is not None and r[gender_col] else "미응답"
        gender_counts[g] += 1
    for g in genders:
        cnt = gender_counts.get(g, 0)
        ws.cell(row=row, column=1, value=g).border = thin_border
        ws.cell(row=row, column=2, value=cnt).border = thin_border
        ws.cell(row=row, column=2).alignment = center_align
        ws.cell(row=row, column=3, value=round(cnt / len(data) * 100, 1)).border = thin_border
        ws.cell(row=row, column=3).alignment = center_align
        row += 1
    row += 1

    # 연령 분포
    ws.cell(row=row, column=1, value="연령 분포").font = Font(bold=True, size=12)
    row += 1
    for c, h in enumerate(["연령", "인원", "비율(%)"], 1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    row += 1

    age_counts = defaultdict(int)
    for r in data:
        a = str(r[age_col]).strip() if age_col is not None and r[age_col] else "미응답"
        age_counts[a] += 1
    for a in ages:
        cnt = age_counts.get(a, 0)
        ws.cell(row=row, column=1, value=a).border = thin_border
        ws.cell(row=row, column=2, value=cnt).border = thin_border
        ws.cell(row=row, column=2).alignment = center_align
        ws.cell(row=row, column=3, value=round(cnt / len(data) * 100, 1)).border = thin_border
        ws.cell(row=row, column=3).alignment = center_align
        row += 1
    row += 1

    # 아이템 목록
    ws.cell(row=row, column=1, value="아이템 목록").font = Font(bold=True, size=12)
    row += 1
    for c, h in enumerate(["No", "첫 질문", "질문 수"], 1):
        cell = ws.cell(row=row, column=c, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = center_align
        cell.border = thin_border
    row += 1

    for item in items:
        ws.cell(row=row, column=1, value=item["item_no"]).border = thin_border
        ws.cell(row=row, column=1).alignment = center_align
        first_q = item["questions"][0]["header"][:70] if item["questions"] else ""
        ws.cell(row=row, column=2, value=first_q).border = thin_border
        ws.cell(row=row, column=3, value=len(item["questions"])).border = thin_border
        ws.cell(row=row, column=3).alignment = center_align
        row += 1

    ws.column_dimensions['A'].width = 15
    ws.column_dimensions['B'].width = 55
    ws.column_dimensions['C'].width = 12


def _write_special_summary(ws, items, data,
                           header_font, header_fill, item_fill, item_font,
                           thin_border, center_align, left_align):
    """특수 질문 요약 시트 — 아이템별 공통 질문(선호도/적정가격/컬러선호) 외 질문을 요약"""
    row = 1
    ws.cell(row=row, column=1, value="아이템별 특수 질문 요약").font = Font(bold=True, size=14, color="1F3864")
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
    row += 2

    for item in items:
        special_qs = [q for q in item["questions"] if q["type"] not in COMMON_TYPES]
        if not special_qs:
            continue

        # 아이템 헤더
        ws.cell(row=row, column=1, value=f"아이템 {item['item_no']}").font = item_font
        ws.cell(row=row, column=1).fill = item_fill
        for c in range(1, 7):
            ws.cell(row=row, column=c).fill = item_fill
            ws.cell(row=row, column=c).border = thin_border
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
        row += 1

        for q in special_qs:
            stats = collect_special_stats(data, q["col_idx"], q["header"], q["type"])
            if stats["type"] == "empty":
                continue

            # 질문 헤더
            ws.cell(row=row, column=1, value=f"[{q['type']}] {q['header'][:100]}").font = Font(bold=True, size=10, color="2F5496")
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=6)
            row += 1

            if stats["type"] == "numeric":
                ws.cell(row=row, column=1, value="전체 평균(0포함)").font = Font(bold=True, size=10)
                ws.cell(row=row, column=1).border = thin_border
                ws.cell(row=row, column=2, value=stats["avg"]).border = thin_border
                ws.cell(row=row, column=2).alignment = center_align
                ws.cell(row=row, column=3, value=f"응답: {stats['count']}명").border = thin_border
                row += 1
                zero_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
                ws.cell(row=row, column=1, value="전체 평균(0제외)").font = Font(bold=True, size=10, color="C00000")
                ws.cell(row=row, column=1).fill = zero_fill
                ws.cell(row=row, column=1).border = thin_border
                ws.cell(row=row, column=2, value=stats["avg_ex"]).border = thin_border
                ws.cell(row=row, column=2).fill = zero_fill
                ws.cell(row=row, column=2).alignment = center_align
                ws.cell(row=row, column=3, value=f"응답: {stats['count_ex']}명").border = thin_border
                ws.cell(row=row, column=3).fill = zero_fill
                row += 1

            elif stats["type"] == "choice":
                # 테이블 헤더
                for ci, h in enumerate(["순위", "응답 항목", "응답수", "비율(%)"], 1):
                    cell = ws.cell(row=row, column=ci, value=h)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = center_align
                    cell.border = thin_border
                row += 1

                for rank, (val, cnt) in enumerate(stats["distribution"][:15], 1):
                    pct = round(cnt / stats["total"] * 100, 1) if stats["total"] > 0 else 0
                    ws.cell(row=row, column=1, value=rank).alignment = center_align
                    ws.cell(row=row, column=1).border = thin_border
                    ws.cell(row=row, column=2, value=val).alignment = left_align
                    ws.cell(row=row, column=2).border = thin_border
                    ws.cell(row=row, column=3, value=cnt).alignment = center_align
                    ws.cell(row=row, column=3).border = thin_border
                    ws.cell(row=row, column=4, value=pct).alignment = center_align
                    ws.cell(row=row, column=4).border = thin_border
                    row += 1

            row += 1  # 질문간 간격

        row += 1  # 아이템간 간격

    ws.column_dimensions['A'].width = 20
    ws.column_dimensions['B'].width = 50
    ws.column_dimensions['C'].width = 12
    ws.column_dimensions['D'].width = 12
    ws.column_dimensions['E'].width = 12
    ws.column_dimensions['F'].width = 12


def _write_dashboard_data(ws, items, data, gender_col, age_col, genders, ages):
    """대시보드용 숨겨진 데이터 시트 작성"""
    # 헤더 구성
    headers = ["item_no", "item_label",
               "선호도_전체", "선호도_응답수"]
    for g in genders:
        headers.append(f"선호도_{g}")
    for a in ages:
        headers.append(f"선호도_{a}")
    for g in genders:
        for a in ages:
            headers.append(f"선호도_{g}_{a}")

    # 0제외 선호도
    headers += ["선호도X_전체", "선호도X_응답수"]
    for g in genders:
        headers.append(f"선호도X_{g}")
    for a in ages:
        headers.append(f"선호도X_{a}")
    for g in genders:
        for a in ages:
            headers.append(f"선호도X_{g}_{a}")

    headers += ["적정가격_전체", "적정가격_응답수"]
    for g in genders:
        headers.append(f"적정가격_{g}")
    for a in ages:
        headers.append(f"적정가격_{a}")
    for g in genders:
        for a in ages:
            headers.append(f"적정가격_{g}_{a}")

    # 0제외 적정가격
    headers += ["적정가격X_전체", "적정가격X_응답수"]
    for g in genders:
        headers.append(f"적정가격X_{g}")
    for a in ages:
        headers.append(f"적정가격X_{a}")
    for g in genders:
        for a in ages:
            headers.append(f"적정가격X_{g}_{a}")

    # 컬러 TOP 8 (이름, 수, 비율)
    for i in range(1, 9):
        headers += [f"컬러{i}_이름", f"컬러{i}_수", f"컬러{i}_비율"]

    # 전체 질문 목록 (최대 6개)
    for i in range(1, 7):
        headers.append(f"질문{i}")

    # 특수질문 슬롯 (최대 4개, 각 2열: 제목+요약)
    for i in range(1, 5):
        headers += [f"특수{i}_제목", f"특수{i}_요약"]

    # 헤더 쓰기
    for c, h in enumerate(headers, 1):
        ws.cell(row=1, column=c, value=h)

    # 데이터 행
    for row_idx, item in enumerate(items, 2):
        col = 1
        ws.cell(row=row_idx, column=col, value=item["item_no"]); col += 1
        ws.cell(row=row_idx, column=col, value=f"아이템{item['item_no']}"); col += 1

        # ── 선호도 0포함 ──
        pref_qs = [q for q in item["questions"] if q["type"] == "선호도"]
        num_block = 2 + len(genders) + len(ages) + len(genders) * len(ages)
        for ez in [False, True]:  # 0포함, 0제외
            if pref_qs:
                avg = calculate_averages(data, pref_qs[0]["col_idx"], gender_col, age_col, exclude_zero=ez)
                ws.cell(row=row_idx, column=col, value=avg["전체"]); col += 1
                ws.cell(row=row_idx, column=col, value=avg["전체_응답수"]); col += 1
                for g in genders:
                    ws.cell(row=row_idx, column=col, value=avg["성별"].get(g, {}).get("평균", "-")); col += 1
                for a in ages:
                    ws.cell(row=row_idx, column=col, value=avg["연령"].get(a, {}).get("평균", "-")); col += 1
                for g in genders:
                    for a in ages:
                        key = f"{g}_{a}"
                        ws.cell(row=row_idx, column=col, value=avg["성별_연령"].get(key, {}).get("평균", "-")); col += 1
            else:
                for _ in range(num_block):
                    ws.cell(row=row_idx, column=col, value="-"); col += 1

        # ── 적정가격 0포함 + 0제외 ──
        price_qs = [q for q in item["questions"] if q["type"] == "적정가격"]
        for ez in [False, True]:  # 0포함, 0제외
            if price_qs:
                avg = calculate_averages(data, price_qs[0]["col_idx"], gender_col, age_col, exclude_zero=ez)
                ws.cell(row=row_idx, column=col, value=avg["전체"]); col += 1
                ws.cell(row=row_idx, column=col, value=avg["전체_응답수"]); col += 1
                for g in genders:
                    ws.cell(row=row_idx, column=col, value=avg["성별"].get(g, {}).get("평균", "-")); col += 1
                for a in ages:
                    ws.cell(row=row_idx, column=col, value=avg["연령"].get(a, {}).get("평균", "-")); col += 1
                for g in genders:
                    for a in ages:
                        key = f"{g}_{a}"
                        ws.cell(row=row_idx, column=col, value=avg["성별_연령"].get(key, {}).get("평균", "-")); col += 1
            else:
                for _ in range(num_block):
                    ws.cell(row=row_idx, column=col, value="-"); col += 1

        # 컬러 TOP 8
        color_qs = [q for q in item["questions"] if q["type"] in ("컬러선호", "컬러선호도")]
        if color_qs:
            stats, total = collect_color_stats(data, color_qs[0]["col_idx"], gender_col, age_col)
            sorted_colors = sorted(stats["전체"].items(), key=lambda x: -x[1])[:8]
            for i in range(8):
                if i < len(sorted_colors):
                    color, count = sorted_colors[i]
                    pct = round(count / total * 100, 1) if total > 0 else 0
                    ws.cell(row=row_idx, column=col, value=color); col += 1
                    ws.cell(row=row_idx, column=col, value=count); col += 1
                    ws.cell(row=row_idx, column=col, value=pct); col += 1
                else:
                    ws.cell(row=row_idx, column=col, value=""); col += 1
                    ws.cell(row=row_idx, column=col, value=""); col += 1
                    ws.cell(row=row_idx, column=col, value=""); col += 1
        else:
            for _ in range(24):
                ws.cell(row=row_idx, column=col, value=""); col += 1

        # 질문 목록 (최대 6개)
        for i in range(6):
            if i < len(item["questions"]):
                ws.cell(row=row_idx, column=col, value=item["questions"][i]["header"][:100])
            col += 1

        # 특수질문 요약 (최대 4개 슬롯: 질문제목 + 요약텍스트)
        special_qs = [q for q in item["questions"] if q["type"] not in COMMON_TYPES]
        for si in range(4):
            if si < len(special_qs):
                sq = special_qs[si]
                stats = collect_special_stats(data, sq["col_idx"], sq["header"], sq["type"])
                ws.cell(row=row_idx, column=col, value=sq["header"][:80]); col += 1
                if stats["type"] == "numeric":
                    ws.cell(row=row_idx, column=col, value=f"평균: {stats['avg']}점 ({stats['count']}명) │ 0제외: {stats['avg_ex']}점 ({stats['count_ex']}명)")
                elif stats["type"] == "choice":
                    top3 = stats["distribution"][:3]
                    parts = []
                    for rank, (val, cnt) in enumerate(top3, 1):
                        pct = round(cnt / stats["total"] * 100, 1) if stats["total"] > 0 else 0
                        parts.append(f"{rank}위: {val[:15]}({pct}%)")
                    ws.cell(row=row_idx, column=col, value=" │ ".join(parts))
                else:
                    ws.cell(row=row_idx, column=col, value="-")
                col += 1
            else:
                ws.cell(row=row_idx, column=col, value=""); col += 1
                ws.cell(row=row_idx, column=col, value=""); col += 1


def _write_dashboard(ws, items, data, gender_col, age_col, genders, ages,
                     header_font, header_fill, item_fill, item_font, sub_header_fill,
                     thin_border, center_align, left_align):
    """아이템 대시보드 시트 - 아이템 번호 선택 시 전체 요약 표시"""

    # DB시트 헤더 위치 맵핑 계산 (0포함 + 0제외 블록)
    num_block = 2 + len(genders) + len(ages) + len(genders) * len(ages)

    db_col = 1  # item_no
    db_col_label = 2  # item_label

    # 0포함 선호도
    db_pref_start = 3  # 선호도_전체
    db_pref_count = 4  # 선호도_응답수
    db_pref_gender_start = 5
    db_pref_age_start = db_pref_gender_start + len(genders)
    db_pref_cross_start = db_pref_age_start + len(ages)

    # 0제외 선호도
    db_pref_ex_start = 3 + num_block
    db_pref_ex_count = db_pref_ex_start + 1
    db_pref_ex_gender_start = db_pref_ex_count + 1
    db_pref_ex_age_start = db_pref_ex_gender_start + len(genders)
    db_pref_ex_cross_start = db_pref_ex_age_start + len(ages)

    # 0포함 적정가격
    db_price_start = 3 + 2 * num_block
    db_price_count = db_price_start + 1
    db_price_gender_start = db_price_count + 1
    db_price_age_start = db_price_gender_start + len(genders)
    db_price_cross_start = db_price_age_start + len(ages)

    # 0제외 적정가격
    db_price_ex_start = 3 + 3 * num_block
    db_price_ex_count = db_price_ex_start + 1
    db_price_ex_gender_start = db_price_ex_count + 1
    db_price_ex_age_start = db_price_ex_gender_start + len(genders)
    db_price_ex_cross_start = db_price_ex_age_start + len(ages)

    db_color_start = 3 + 4 * num_block
    db_question_start = db_color_start + 24  # 8 colors * 3 columns
    db_special_start = db_question_start + 6  # 6 questions, then special question slots

    # 스타일
    title_font = Font(bold=True, size=16, color="1F3864")
    section_font = Font(bold=True, size=13, color="FFFFFF")
    section_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    label_font = Font(bold=True, size=11, color="333333")
    value_font = Font(size=11)
    input_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
    input_border = Border(
        left=Side(style='medium', color='D4A017'),
        right=Side(style='medium', color='D4A017'),
        top=Side(style='medium', color='D4A017'),
        bottom=Side(style='medium', color='D4A017')
    )
    accent_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    light_blue_fill = PatternFill(start_color="DAEEF3", end_color="DAEEF3", fill_type="solid")
    light_gray_fill = PatternFill(start_color="F2F2F2", end_color="F2F2F2", fill_type="solid")

    db_sheet_name = "_DB"
    total_items = len(items)
    # VLOOKUP helper
    def vlookup(db_col_num):
        return f'IFERROR(VLOOKUP($C$3,{db_sheet_name}!$A:${get_column_letter(db_special_start + 8)},{db_col_num},FALSE),"-")'

    row = 1
    # ===================== 타이틀 =====================
    ws.cell(row=row, column=1, value="📊 아이템 대시보드").font = title_font
    ws.merge_cells("A1:H1")
    row += 2

    # ===================== 아이템 선택 영역 =====================
    ws.cell(row=row, column=1, value="▶ 아이템 번호 선택:").font = Font(bold=True, size=12, color="1F3864")
    ws.cell(row=row, column=1).alignment = Alignment(horizontal='right', vertical='center')

    input_cell = ws.cell(row=row, column=3, value=1)
    input_cell.font = Font(bold=True, size=14, color="C00000")
    input_cell.fill = input_fill
    input_cell.border = input_border
    input_cell.alignment = center_align

    # 드롭다운 데이터 검증
    item_list = ",".join([str(i+1) for i in range(total_items)])
    dv = DataValidation(type="list", formula1=f'"{item_list}"', allow_blank=False)
    dv.error = "1~{}사이 아이템 번호를 선택해주세요".format(total_items)
    dv.errorTitle = "잘못된 입력"
    dv.prompt = "분석할 아이템 번호를 선택하세요"
    dv.promptTitle = "아이템 선택"
    ws.add_data_validation(dv)
    dv.add(ws["C3"])

    # 아이템 라벨 표시
    ws.cell(row=row, column=4, value=f'={vlookup(db_col_label)}').font = Font(bold=True, size=12, color="2F5496")
    ws.cell(row=row, column=5, value=f"(1~{total_items} 선택 가능)").font = Font(italic=True, size=10, color="808080")
    row += 2

    # ===================== 섹션 1: 선호도 =====================
    ws.cell(row=row, column=1, value="■ 선호도 평균 (10점 만점)").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = section_fill
        ws.cell(row=row, column=c).font = section_font
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    ex_yellow = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
    ex_label_font = Font(bold=True, size=9, color="C00000")

    # 전체 평균 (0포함)
    ws.cell(row=row, column=1, value="전체(0포함)").font = label_font
    ws.cell(row=row, column=1).fill = light_gray_fill
    ws.cell(row=row, column=1).border = thin_border
    ws.cell(row=row, column=2).border = thin_border
    ws.cell(row=row, column=2, value=f'={vlookup(db_pref_start)}').font = Font(bold=True, size=14, color="C00000")
    ws.cell(row=row, column=2).alignment = center_align
    ws.cell(row=row, column=2).fill = light_gray_fill
    ws.cell(row=row, column=3, value="응답수").font = label_font
    ws.cell(row=row, column=3).fill = light_gray_fill
    ws.cell(row=row, column=3).border = thin_border
    ws.cell(row=row, column=4, value=f'={vlookup(db_pref_count)}').alignment = center_align
    ws.cell(row=row, column=4).fill = light_gray_fill
    ws.cell(row=row, column=4).border = thin_border
    row += 1

    # 전체 평균 (0제외)
    ws.cell(row=row, column=1, value="전체(0제외)").font = ex_label_font
    ws.cell(row=row, column=1).fill = ex_yellow
    ws.cell(row=row, column=1).border = thin_border
    ws.cell(row=row, column=2, value=f'={vlookup(db_pref_ex_start)}').font = Font(bold=True, size=14, color="C00000")
    ws.cell(row=row, column=2).alignment = center_align
    ws.cell(row=row, column=2).fill = ex_yellow
    ws.cell(row=row, column=2).border = thin_border
    ws.cell(row=row, column=3, value="응답수").font = ex_label_font
    ws.cell(row=row, column=3).fill = ex_yellow
    ws.cell(row=row, column=3).border = thin_border
    ws.cell(row=row, column=4, value=f'={vlookup(db_pref_ex_count)}').alignment = center_align
    ws.cell(row=row, column=4).fill = ex_yellow
    ws.cell(row=row, column=4).border = thin_border
    row += 1

    # 성별 평균 (0포함)
    ws.cell(row=row, column=1, value="성별(0포함)").font = label_font
    ws.cell(row=row, column=1).fill = accent_fill
    ws.cell(row=row, column=1).border = thin_border
    for gi, g in enumerate(genders):
        col_base = 2 + gi * 2
        ws.cell(row=row, column=col_base, value=g).font = label_font
        ws.cell(row=row, column=col_base).alignment = center_align
        ws.cell(row=row, column=col_base).fill = accent_fill
        ws.cell(row=row, column=col_base).border = thin_border
        ws.cell(row=row, column=col_base + 1, value=f'={vlookup(db_pref_gender_start + gi)}').alignment = center_align
        ws.cell(row=row, column=col_base + 1).font = Font(bold=True, size=12)
        ws.cell(row=row, column=col_base + 1).fill = accent_fill
        ws.cell(row=row, column=col_base + 1).border = thin_border
    row += 1

    # 성별 평균 (0제외)
    ws.cell(row=row, column=1, value="성별(0제외)").font = ex_label_font
    ws.cell(row=row, column=1).fill = ex_yellow
    ws.cell(row=row, column=1).border = thin_border
    for gi, g in enumerate(genders):
        col_base = 2 + gi * 2
        ws.cell(row=row, column=col_base, value=g).font = Font(bold=True, size=9, color="C00000")
        ws.cell(row=row, column=col_base).alignment = center_align
        ws.cell(row=row, column=col_base).fill = ex_yellow
        ws.cell(row=row, column=col_base).border = thin_border
        ws.cell(row=row, column=col_base + 1, value=f'={vlookup(db_pref_ex_gender_start + gi)}').alignment = center_align
        ws.cell(row=row, column=col_base + 1).font = Font(bold=True, size=12, color="C00000")
        ws.cell(row=row, column=col_base + 1).fill = ex_yellow
        ws.cell(row=row, column=col_base + 1).border = thin_border
    row += 1

    # 연령별 평균 헤더
    ws.cell(row=row, column=1, value="연령(0포함)").font = label_font
    ws.cell(row=row, column=1).fill = light_blue_fill
    ws.cell(row=row, column=1).border = thin_border
    for ai, a in enumerate(ages):
        c = 2 + ai
        cell = ws.cell(row=row, column=c, value=a)
        cell.font = Font(size=9, bold=True)
        cell.alignment = center_align
        cell.fill = light_blue_fill
        cell.border = thin_border
    row += 1
    ws.cell(row=row, column=1).fill = light_blue_fill
    ws.cell(row=row, column=1).border = thin_border
    for ai in range(len(ages)):
        c = 2 + ai
        cell = ws.cell(row=row, column=c, value=f'={vlookup(db_pref_age_start + ai)}')
        cell.font = Font(bold=True, size=11)
        cell.alignment = center_align
        cell.fill = light_blue_fill
        cell.border = thin_border
    row += 1

    # 연령별 평균 (0제외)
    ws.cell(row=row, column=1, value="연령(0제외)").font = ex_label_font
    ws.cell(row=row, column=1).fill = ex_yellow
    ws.cell(row=row, column=1).border = thin_border
    for ai in range(len(ages)):
        c = 2 + ai
        cell = ws.cell(row=row, column=c, value=f'={vlookup(db_pref_ex_age_start + ai)}')
        cell.font = Font(bold=True, size=11, color="C00000")
        cell.alignment = center_align
        cell.fill = ex_yellow
        cell.border = thin_border
    row += 2

    # ===================== 섹션 2: 적정가격 =====================
    ws.cell(row=row, column=1, value="■ 적정가격 평균 (원)").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = section_fill
        ws.cell(row=row, column=c).font = section_font
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    # 전체 평균 (0포함)
    ws.cell(row=row, column=1, value="전체(0포함)").font = label_font
    ws.cell(row=row, column=1).fill = light_gray_fill
    ws.cell(row=row, column=1).border = thin_border
    price_cell = ws.cell(row=row, column=2, value=f'={vlookup(db_price_start)}')
    price_cell.font = Font(bold=True, size=14, color="C00000")
    price_cell.alignment = center_align
    price_cell.fill = light_gray_fill
    price_cell.border = thin_border
    price_cell.number_format = '#,##0'
    ws.cell(row=row, column=3, value="응답수").font = label_font
    ws.cell(row=row, column=3).fill = light_gray_fill
    ws.cell(row=row, column=3).border = thin_border
    ws.cell(row=row, column=4, value=f'={vlookup(db_price_count)}').alignment = center_align
    ws.cell(row=row, column=4).fill = light_gray_fill
    ws.cell(row=row, column=4).border = thin_border
    row += 1

    # 전체 평균 (0제외)
    ws.cell(row=row, column=1, value="전체(0제외)").font = ex_label_font
    ws.cell(row=row, column=1).fill = ex_yellow
    ws.cell(row=row, column=1).border = thin_border
    price_ex_cell = ws.cell(row=row, column=2, value=f'={vlookup(db_price_ex_start)}')
    price_ex_cell.font = Font(bold=True, size=14, color="C00000")
    price_ex_cell.alignment = center_align
    price_ex_cell.fill = ex_yellow
    price_ex_cell.border = thin_border
    price_ex_cell.number_format = '#,##0'
    ws.cell(row=row, column=3, value="응답수").font = ex_label_font
    ws.cell(row=row, column=3).fill = ex_yellow
    ws.cell(row=row, column=3).border = thin_border
    ws.cell(row=row, column=4, value=f'={vlookup(db_price_ex_count)}').alignment = center_align
    ws.cell(row=row, column=4).fill = ex_yellow
    ws.cell(row=row, column=4).border = thin_border
    row += 1

    # 성별 (0포함)
    ws.cell(row=row, column=1, value="성별(0포함)").font = label_font
    ws.cell(row=row, column=1).fill = accent_fill
    ws.cell(row=row, column=1).border = thin_border
    for gi, g in enumerate(genders):
        col_base = 2 + gi * 2
        ws.cell(row=row, column=col_base, value=g).font = label_font
        ws.cell(row=row, column=col_base).alignment = center_align
        ws.cell(row=row, column=col_base).fill = accent_fill
        ws.cell(row=row, column=col_base).border = thin_border
        c = ws.cell(row=row, column=col_base + 1, value=f'={vlookup(db_price_gender_start + gi)}')
        c.alignment = center_align
        c.font = Font(bold=True, size=12)
        c.fill = accent_fill
        c.border = thin_border
        c.number_format = '#,##0'
    row += 1

    # 성별 (0제외)
    ws.cell(row=row, column=1, value="성별(0제외)").font = ex_label_font
    ws.cell(row=row, column=1).fill = ex_yellow
    ws.cell(row=row, column=1).border = thin_border
    for gi, g in enumerate(genders):
        col_base = 2 + gi * 2
        ws.cell(row=row, column=col_base, value=g).font = Font(bold=True, size=9, color="C00000")
        ws.cell(row=row, column=col_base).alignment = center_align
        ws.cell(row=row, column=col_base).fill = ex_yellow
        ws.cell(row=row, column=col_base).border = thin_border
        c = ws.cell(row=row, column=col_base + 1, value=f'={vlookup(db_price_ex_gender_start + gi)}')
        c.alignment = center_align
        c.font = Font(bold=True, size=12, color="C00000")
        c.fill = ex_yellow
        c.border = thin_border
        c.number_format = '#,##0'
    row += 1

    # 연령별 (0포함)
    ws.cell(row=row, column=1, value="연령(0포함)").font = label_font
    ws.cell(row=row, column=1).fill = light_blue_fill
    ws.cell(row=row, column=1).border = thin_border
    for ai, a in enumerate(ages):
        c = 2 + ai
        cell = ws.cell(row=row, column=c, value=a)
        cell.font = Font(size=9, bold=True)
        cell.alignment = center_align
        cell.fill = light_blue_fill
        cell.border = thin_border
    row += 1
    ws.cell(row=row, column=1).fill = light_blue_fill
    ws.cell(row=row, column=1).border = thin_border
    for ai in range(len(ages)):
        c = 2 + ai
        cell = ws.cell(row=row, column=c, value=f'={vlookup(db_price_age_start + ai)}')
        cell.font = Font(bold=True, size=11)
        cell.alignment = center_align
        cell.fill = light_blue_fill
        cell.border = thin_border
        cell.number_format = '#,##0'
    row += 1

    # 연령별 (0제외)
    ws.cell(row=row, column=1, value="연령(0제외)").font = ex_label_font
    ws.cell(row=row, column=1).fill = ex_yellow
    ws.cell(row=row, column=1).border = thin_border
    for ai in range(len(ages)):
        c = 2 + ai
        cell = ws.cell(row=row, column=c, value=f'={vlookup(db_price_ex_age_start + ai)}')
        cell.font = Font(bold=True, size=11, color="C00000")
        cell.alignment = center_align
        cell.fill = ex_yellow
        cell.border = thin_border
        cell.number_format = '#,##0'
    row += 2

    # ===================== 섹션 3: 컬러 선호 TOP 8 =====================
    ws.cell(row=row, column=1, value="■ 컬러 선호 TOP 8").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = section_fill
        ws.cell(row=row, column=c).font = section_font
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    # 헤더
    color_headers = ["순위", "컬러명", "응답수", "비율(%)"] * 2
    for ci, h in enumerate(color_headers, 1):
        cell = ws.cell(row=row, column=ci, value=h)
        cell.font = Font(bold=True, size=10, color="FFFFFF")
        cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        cell.alignment = center_align
        cell.border = thin_border
    row += 1

    # TOP 8 컬러 (좌 1~4, 우 5~8)
    for i in range(4):
        # 좌측 1~4
        ws.cell(row=row, column=1, value=i+1).alignment = center_align
        ws.cell(row=row, column=1).border = thin_border
        ws.cell(row=row, column=1).font = Font(bold=True)
        db_c = db_color_start + i * 3
        ws.cell(row=row, column=2, value=f'={vlookup(db_c)}').alignment = center_align
        ws.cell(row=row, column=2).border = thin_border
        ws.cell(row=row, column=3, value=f'={vlookup(db_c + 1)}').alignment = center_align
        ws.cell(row=row, column=3).border = thin_border
        cell_pct = ws.cell(row=row, column=4, value=f'={vlookup(db_c + 2)}')
        cell_pct.alignment = center_align
        cell_pct.border = thin_border

        # 우측 5~8
        ws.cell(row=row, column=5, value=i+5).alignment = center_align
        ws.cell(row=row, column=5).border = thin_border
        ws.cell(row=row, column=5).font = Font(bold=True)
        db_c2 = db_color_start + (i + 4) * 3
        ws.cell(row=row, column=6, value=f'={vlookup(db_c2)}').alignment = center_align
        ws.cell(row=row, column=6).border = thin_border
        ws.cell(row=row, column=7, value=f'={vlookup(db_c2 + 1)}').alignment = center_align
        ws.cell(row=row, column=7).border = thin_border
        cell_pct2 = ws.cell(row=row, column=8, value=f'={vlookup(db_c2 + 2)}')
        cell_pct2.alignment = center_align
        cell_pct2.border = thin_border

        # 줄무늬
        if i % 2 == 1:
            for cc in range(1, 9):
                ws.cell(row=row, column=cc).fill = light_gray_fill
        row += 1
    row += 1

    # ===================== 섹션 4: 성별×연령 교차분석 =====================
    ws.cell(row=row, column=1, value="■ 성별 × 연령 교차분석 - 선호도 (0포함)").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = section_fill
        ws.cell(row=row, column=c).font = section_font
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    # 교차분석 테이블 (0포함)
    cross_col = 2
    ws.cell(row=row, column=1, value="").border = thin_border
    ws.cell(row=row, column=1).fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    for a in ages:
        cell = ws.cell(row=row, column=cross_col, value=a)
        cell.font = Font(bold=True, size=9, color="FFFFFF")
        cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        cell.alignment = center_align
        cell.border = thin_border
        cross_col += 1
    row += 1

    for gi, g in enumerate(genders):
        ws.cell(row=row, column=1, value=g).font = Font(bold=True, size=10)
        ws.cell(row=row, column=1).alignment = center_align
        ws.cell(row=row, column=1).border = thin_border
        if gi % 2 == 0:
            ws.cell(row=row, column=1).fill = accent_fill
        else:
            ws.cell(row=row, column=1).fill = light_blue_fill

        for ai in range(len(ages)):
            db_c = db_pref_cross_start + gi * len(ages) + ai
            cell = ws.cell(row=row, column=2 + ai, value=f'={vlookup(db_c)}')
            cell.alignment = center_align
            cell.border = thin_border
            cell.font = Font(bold=True, size=11)
            if gi % 2 == 0:
                cell.fill = accent_fill
            else:
                cell.fill = light_blue_fill
        row += 1
    row += 1

    # 교차분석 테이블 (0제외)
    ws.cell(row=row, column=1, value="■ 성별 × 연령 교차분석 - 선호도 (0제외)").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = PatternFill(start_color="BF8F00", end_color="BF8F00", fill_type="solid")
        ws.cell(row=row, column=c).font = Font(bold=True, size=13, color="FFFFFF")
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    cross_col = 2
    ws.cell(row=row, column=1, value="").border = thin_border
    ws.cell(row=row, column=1).fill = PatternFill(start_color="BF8F00", end_color="BF8F00", fill_type="solid")
    for a in ages:
        cell = ws.cell(row=row, column=cross_col, value=a)
        cell.font = Font(bold=True, size=9, color="FFFFFF")
        cell.fill = PatternFill(start_color="BF8F00", end_color="BF8F00", fill_type="solid")
        cell.alignment = center_align
        cell.border = thin_border
        cross_col += 1
    row += 1

    for gi, g in enumerate(genders):
        ws.cell(row=row, column=1, value=g).font = Font(bold=True, size=10, color="C00000")
        ws.cell(row=row, column=1).alignment = center_align
        ws.cell(row=row, column=1).border = thin_border
        ws.cell(row=row, column=1).fill = ex_yellow

        for ai in range(len(ages)):
            db_c = db_pref_ex_cross_start + gi * len(ages) + ai
            cell = ws.cell(row=row, column=2 + ai, value=f'={vlookup(db_c)}')
            cell.alignment = center_align
            cell.border = thin_border
            cell.font = Font(bold=True, size=11, color="C00000")
            cell.fill = ex_yellow
        row += 1
    row += 1

    # ===================== 섹션 5: 해당 아이템 질문 목록 =====================
    ws.cell(row=row, column=1, value="■ 해당 아이템 질문 목록").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = section_fill
        ws.cell(row=row, column=c).font = section_font
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    for i in range(6):
        ws.cell(row=row, column=1, value=f"질문 {i+1}").font = label_font
        ws.cell(row=row, column=1).alignment = Alignment(horizontal='center', vertical='center')
        ws.cell(row=row, column=1).border = thin_border
        ws.cell(row=row, column=1).fill = light_gray_fill
        q_cell = ws.cell(row=row, column=2, value=f'={vlookup(db_question_start + i)}')
        q_cell.alignment = left_align
        q_cell.border = thin_border
        ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=8)
        for cc in range(2, 9):
            ws.cell(row=row, column=cc).border = thin_border
        row += 1

    # ===================== 섹션 6: 특수질문 요약 =====================
    ws.cell(row=row, column=1, value="■ 특수질문 요약 (해당 아이템)").font = section_font
    for c in range(1, 9):
        ws.cell(row=row, column=c).fill = PatternFill(start_color="548235", end_color="548235", fill_type="solid")
        ws.cell(row=row, column=c).font = Font(bold=True, size=13, color="FFFFFF")
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=8)
    row += 1

    special_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    for si in range(4):
        db_title_col = db_special_start + si * 2
        db_summary_col = db_special_start + si * 2 + 1
        # 질문 제목
        ws.cell(row=row, column=1, value=f"특수질문 {si+1}").font = Font(bold=True, size=10, color="548235")
        ws.cell(row=row, column=1).fill = special_fill
        ws.cell(row=row, column=1).border = thin_border
        ws.cell(row=row, column=1).alignment = Alignment(horizontal='center', vertical='center')
        title_cell = ws.cell(row=row, column=2, value=f'={vlookup(db_title_col)}')
        title_cell.font = Font(bold=True, size=10)
        title_cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
        title_cell.border = thin_border
        ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=8)
        for cc in range(2, 9):
            ws.cell(row=row, column=cc).border = thin_border
        row += 1
        # 요약
        ws.cell(row=row, column=1).fill = special_fill
        ws.cell(row=row, column=1).border = thin_border
        summary_cell = ws.cell(row=row, column=2, value=f'={vlookup(db_summary_col)}')
        summary_cell.font = Font(size=10, color="333333")
        summary_cell.alignment = Alignment(horizontal='left', vertical='center')
        summary_cell.border = thin_border
        ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=8)
        for cc in range(2, 9):
            ws.cell(row=row, column=cc).border = thin_border
        row += 1

    # 열 너비
    ws.column_dimensions['A'].width = 14
    ws.column_dimensions['B'].width = 16
    ws.column_dimensions['C'].width = 14
    ws.column_dimensions['D'].width = 14
    ws.column_dimensions['E'].width = 14
    ws.column_dimensions['F'].width = 14
    ws.column_dimensions['G'].width = 14
    ws.column_dimensions['H'].width = 14

    # 행 높이
    ws.row_dimensions[1].height = 30
    ws.row_dimensions[3].height = 28

    # 인쇄 설정
    ws.sheet_properties.pageSetUpPr = openpyxl.worksheet.properties.PageSetupProperties(fitToPage=True)


# ================================================================
# PPT 생성
# ================================================================

def _add_ranking_table(slide, left_in, top_in, width_in, height_in,
                       title, title_bg_rgb, items, ages, score_min, score_max):
    """요약 페이지 랭킹 테이블 (데이터 바 포함)"""
    n_items = len(items)
    n_age_cols = len(ages)
    n_cols = 3 + n_age_cols

    rank_w = 0.38
    sty_w = 1.22
    ttl_w = 0.50
    remain = width_in - rank_w - sty_w - ttl_w
    age_w = remain / max(n_age_cols, 1)

    header_row_h = Inches(0.22)
    data_row_h = Inches(0.17)
    n_data = n_items + 1
    n_rows = 1 + n_data
    tbl_h = header_row_h + data_row_h * n_data

    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)

    hdr_h = Inches(0.28)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = title_bg_rgb
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_PPT_WHITE

    tbl_top = top + hdr_h
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    tbl = tbl_shape.table
    tbl.first_row = False
    _clear_table_style(tbl_shape)

    tbl.rows[0].height = header_row_h
    for r in range(1, n_rows):
        tbl.rows[r].height = data_row_h

    tbl.columns[0].width = Inches(rank_w)
    tbl.columns[1].width = Inches(sty_w)
    tbl.columns[2].width = Inches(ttl_w)
    for i in range(3, n_cols):
        tbl.columns[i].width = Inches(age_w)

    header_texts = ["순위", "STY", "TTL"] + [a.replace("세", "") for a in ages]
    for ci, txt in enumerate(header_texts):
        c = tbl.cell(0, ci)
        c.text = txt
        _style_ppt_cell(c, Pt(7), True, C_PPT_WHITE, title_bg_rgb)
        _set_cell_border(c, "000000", "6350")

    for ri, item in enumerate(items):
        row_idx = ri + 1
        rank = ri + 1

        c_rank = tbl.cell(row_idx, 0)
        c_rank.text = str(rank)
        _style_ppt_cell(c_rank, Pt(7), True, C_PPT_BLACK, C_PPT_WHITE)
        _set_cell_border(c_rank, "000000", "6350")

        c_name = tbl.cell(row_idx, 1)
        c_name.text = item["name"]
        _style_ppt_cell(c_name, Pt(7), False, C_PPT_BLACK, C_PPT_WHITE, align=PP_ALIGN.LEFT)
        c_name.text_frame.margin_left = Pt(4)
        _set_cell_border(c_name, "000000", "6350")

        ttl_val = item["ttl"]
        c_ttl = tbl.cell(row_idx, 2)
        c_ttl.text = str(ttl_val) if isinstance(ttl_val, (int, float)) else "-"
        is_high_ttl = isinstance(ttl_val, (int, float)) and ttl_val >= 5
        _style_ppt_cell(c_ttl, Pt(7), is_high_ttl, C_SCORE_RED if is_high_ttl else C_PPT_BLACK, C_PPT_WHITE)
        _set_cell_border(c_ttl, "000000", "6350")
        ratio_ttl = _data_bar_width_ratio(ttl_val, score_min, score_max)
        _apply_cell_data_bar(c_ttl, ratio_ttl, C_BAR_GOLD, C_PPT_WHITE)

        for ai, a in enumerate(ages):
            ci = 3 + ai
            c = tbl.cell(row_idx, ci)
            val = item["ages"].get(a, "-")
            c.text = str(val) if isinstance(val, (int, float)) else "-"
            is_high_v = isinstance(val, (int, float)) and val >= 5
            _style_ppt_cell(c, Pt(7), is_high_v, C_SCORE_RED if is_high_v else C_PPT_BLACK, C_PPT_WHITE)
            _set_cell_border(c, "000000", "6350")
            ratio_v = _data_bar_width_ratio(val, score_min, score_max)
            _apply_cell_data_bar(c, ratio_v, C_BAR_BLUE, C_PPT_WHITE)

    n_row_idx = n_items + 1
    cn0 = tbl.cell(n_row_idx, 0)
    cn0.text = ""
    _style_ppt_cell(cn0, Pt(6), False, C_PPT_BLACK, C_LIGHT_GRAY)
    _set_cell_border(cn0, "000000", "6350")
    cn1 = tbl.cell(n_row_idx, 1)
    cn1.text = "인원수 (n)"
    _style_ppt_cell(cn1, Pt(6), True, C_PPT_BLACK, C_LIGHT_GRAY, align=PP_ALIGN.LEFT)
    cn1.text_frame.margin_left = Pt(4)
    _set_cell_border(cn1, "000000", "6350")
    cn_ttl = tbl.cell(n_row_idx, 2)
    n_total = items[0]["n_total"] if items else 0
    cn_ttl.text = str(n_total)
    _style_ppt_cell(cn_ttl, Pt(6), True, C_PPT_BLACK, C_LIGHT_GRAY)
    _set_cell_border(cn_ttl, "000000", "6350")
    for ai, a in enumerate(ages):
        ci = 3 + ai
        cn = tbl.cell(n_row_idx, ci)
        n_val = items[0]["n_ages"].get(a, 0) if items else 0
        cn.text = str(n_val)
        _style_ppt_cell(cn, Pt(6), False, C_PPT_BLACK, C_LIGHT_GRAY)
        _set_cell_border(cn, "000000", "6350")


def _add_item_data_table(slide, left_in, top_in, width_in, title, title_bg,
                         avg_data, avg_data_ex, genders, ages, fmt="score"):
    """선호도/적정가격 통합 테이블"""
    n_cols = 1 + 2 + len(genders) * 2
    row_labels = ["전체"] + ages
    n_rows = 1 + len(row_labels)

    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    row_h = Inches(0.19)

    hdr_h = Inches(0.25)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = title_bg
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_PPT_WHITE

    tbl_top = top + hdr_h
    tbl_h = row_h * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    tbl = tbl_shape.table
    tbl.first_row = False
    _clear_table_style(tbl_shape)

    for i in range(n_rows):
        tbl.rows[i].height = row_h

    label_w = 0.50
    remaining = width_in - label_w
    col_w = remaining / (n_cols - 1)
    tbl.columns[0].width = Inches(label_w)
    for c in range(1, n_cols):
        tbl.columns[c].width = Inches(col_w)

    hdrs = ["구분", "전체", "(0제외)"]
    for g in genders:
        hdrs.append(g)
        hdrs.append("(0제외)")
    for ci, txt in enumerate(hdrs):
        c = tbl.cell(0, ci)
        c.text = txt
        is_exc = (ci == 2) or (ci >= 3 and ci % 2 == 0)
        bg = C_HDR_EXC if is_exc else C_HDR_MAIN
        _style_ppt_cell(c, Pt(7), True, C_PPT_WHITE, bg)
        _set_cell_border(c, "000000", "6350", dash=None)

    def _fmt(val):
        if isinstance(val, (int, float)):
            return f"{int(val):,}" if fmt == "price" else str(val)
        return str(val)

    col_vals = [[] for _ in range(n_cols - 1)]
    for label in row_labels:
        if label == "전체":
            vals = [avg_data["전체"], avg_data_ex["전체"]]
            for g in genders:
                vals.append(avg_data["성별"].get(g, {}).get("평균", "-"))
                vals.append(avg_data_ex["성별"].get(g, {}).get("평균", "-"))
        else:
            a = label
            vals = [
                avg_data["연령"].get(a, {}).get("평균", "-"),
                avg_data_ex["연령"].get(a, {}).get("평균", "-"),
            ]
            for g in genders:
                key = f"{g}_{a}"
                vals.append(avg_data["성별_연령"].get(key, {}).get("평균", "-"))
                vals.append(avg_data_ex["성별_연령"].get(key, {}).get("평균", "-"))
        for ci, v in enumerate(vals):
            if isinstance(v, (int, float)):
                col_vals[ci].append(v)
    col_maxes = [max(vs) if vs else None for vs in col_vals]

    for ri, label in enumerate(row_labels):
        row_idx = ri + 1
        is_total = (label == "전체")

        c0 = tbl.cell(row_idx, 0)
        c0.text = label if is_total else label.replace("세", "")
        label_bg = C_TOTAL_LABEL if is_total else C_PPT_WHITE
        _style_ppt_cell(c0, Pt(7), True, C_NAVY, label_bg)
        _set_cell_border(c0, "000000", "6350")

        if is_total:
            vals = [avg_data["전체"], avg_data_ex["전체"]]
            for g in genders:
                vals.append(avg_data["성별"].get(g, {}).get("평균", "-"))
                vals.append(avg_data_ex["성별"].get(g, {}).get("평균", "-"))
        else:
            a = label
            vals = [
                avg_data["연령"].get(a, {}).get("평균", "-"),
                avg_data_ex["연령"].get(a, {}).get("평균", "-"),
            ]
            for g in genders:
                key = f"{g}_{a}"
                vals.append(avg_data["성별_연령"].get(key, {}).get("평균", "-"))
                vals.append(avg_data_ex["성별_연령"].get(key, {}).get("평균", "-"))

        for ci, v in enumerate(vals):
            col_idx = ci + 1
            c = tbl.cell(row_idx, col_idx)
            c.text = _fmt(v)
            is_max = isinstance(v, (int, float)) and col_maxes[ci] is not None and v == col_maxes[ci]
            fc = C_SCORE_RED if (is_total or is_max) else C_PPT_BLACK
            bold = is_total or is_max
            bg = C_TOTAL_BG if is_total else C_PPT_WHITE
            _style_ppt_cell(c, Pt(7), bold, fc, bg)
            _set_cell_border(c, "000000", "6350")

    return Inches(top_in) + hdr_h + tbl_h


def _add_color_table(slide, left_in, top_in, width_in, title_bg,
                     color_stats, color_total):
    """컬러 선호 TOP 10 테이블"""
    sorted_colors = sorted(color_stats["전체"].items(), key=lambda x: -x[1])[:10]
    if not sorted_colors:
        return Inches(top_in)

    n_rows = len(sorted_colors) + 1
    n_cols = 3
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    row_h = Inches(0.19)

    hdr_h = Inches(0.25)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = title_bg
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "컬러 선호 TOP"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_PPT_WHITE

    tbl_top = top + hdr_h
    tbl_h = row_h * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, tbl_top, width, tbl_h)
    tbl = tbl_shape.table
    tbl.first_row = False
    _clear_table_style(tbl_shape)

    for i in range(n_rows):
        tbl.rows[i].height = row_h

    tbl.columns[0].width = Inches(width_in * 0.50)
    tbl.columns[1].width = Inches(width_in * 0.25)
    tbl.columns[2].width = Inches(width_in * 0.25)

    for ci, txt in enumerate(["컬러명", "응답수", "비율"]):
        c = tbl.cell(0, ci)
        c.text = txt
        _style_ppt_cell(c, Pt(7.5), True, C_PPT_WHITE, C_HDR_MAIN)
        _set_cell_border(c, "000000", "6350", dash=None)

    max_count = sorted_colors[0][1] if sorted_colors else 0
    for ri, (color_name, cnt) in enumerate(sorted_colors):
        row_idx = ri + 1
        pct = round(cnt / color_total * 100, 1) if color_total > 0 else 0

        c0 = tbl.cell(row_idx, 0)
        c0.text = f"{ri+1}. {color_name}"
        _style_ppt_cell(c0, Pt(7), False, C_PPT_BLACK, C_PPT_WHITE, align=PP_ALIGN.LEFT)
        c0.text_frame.margin_left = Pt(4)
        _set_cell_border(c0, "000000", "6350")

        c1 = tbl.cell(row_idx, 1)
        c1.text = f"{cnt}명"
        is_max = (cnt == max_count)
        _style_ppt_cell(c1, Pt(7), is_max, C_SCORE_RED if is_max else C_PPT_BLACK, C_PPT_WHITE)
        _set_cell_border(c1, "000000", "6350")

        c2 = tbl.cell(row_idx, 2)
        c2.text = f"{pct}%"
        _style_ppt_cell(c2, Pt(7), is_max, C_SCORE_RED if is_max else C_PPT_BLACK, C_PPT_WHITE)
        _set_cell_border(c2, "000000", "6350")

    return tbl_top + tbl_h


def _add_cross_table(slide, left_in, top_in, width_in, pref_avg, genders, ages, title_bg):
    """성별×연령 교차분석 테이블 (성별별 최고값만 하이라이트)"""
    rows = len(genders) + 1
    cols = len(ages) + 1
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    row_h = Inches(0.22)

    hdr_h = Inches(0.25)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, hdr_h)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = title_bg
    hdr.line.fill.background()
    tf = hdr.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "성별×연령 교차분석 (0포함)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_PPT_WHITE

    tbl_top = top + hdr_h
    tbl_h = row_h * rows
    tbl_shape = slide.shapes.add_table(rows, cols, left, tbl_top, width, tbl_h)
    tbl = tbl_shape.table
    tbl.first_row = False
    _clear_table_style(tbl_shape)

    for i in range(rows):
        tbl.rows[i].height = row_h

    tbl.columns[0].width = Inches(0.50)
    age_col_w = Inches((width_in - 0.50) / max(len(ages), 1))
    for c in range(1, cols):
        tbl.columns[c].width = age_col_w

    c0 = tbl.cell(0, 0)
    c0.text = "성별\\연령"
    _style_ppt_cell(c0, Pt(7), True, C_PPT_WHITE, C_HDR_MAIN)
    _set_cell_border(c0, "000000", "6350", dash=None)
    for ai, a in enumerate(ages):
        c = tbl.cell(0, ai + 1)
        c.text = a.replace("세", "")
        _style_ppt_cell(c, Pt(7), True, C_PPT_WHITE, C_HDR_MAIN)
        _set_cell_border(c, "000000", "6350", dash=None)

    gender_max = {}
    for gi, g in enumerate(genders):
        max_v = None
        for a in ages:
            key = f"{g}_{a}"
            val = pref_avg["성별_연령"].get(key, {}).get("평균", "-")
            if isinstance(val, (int, float)):
                if max_v is None or val > max_v:
                    max_v = val
        gender_max[g] = max_v

    for gi, g in enumerate(genders):
        ri = gi + 1
        c0 = tbl.cell(ri, 0)
        c0.text = g[:4] if len(g) > 6 else g
        _style_ppt_cell(c0, Pt(7), True, C_NAVY, C_PPT_WHITE)
        _set_cell_border(c0, "000000", "6350")

        for ai, a in enumerate(ages):
            key = f"{g}_{a}"
            val = pref_avg["성별_연령"].get(key, {}).get("평균", "-")
            c = tbl.cell(ri, ai + 1)
            c.text = str(val)
            is_max = isinstance(val, (int, float)) and gender_max[g] is not None and val == gender_max[g]
            if is_max:
                bg = RGBColor(0xC6, 0xEF, 0xCE)
                fc = C_SCORE_RED
            else:
                bg = C_PPT_WHITE
                fc = C_PPT_BLACK
            _style_ppt_cell(c, Pt(8), is_max, fc, bg)
            _set_cell_border(c, "000000", "6350")

    return tbl_top + tbl_h


def _add_item_slide(prs, slide, item, item_no, data, gender_col, age_col,
                    genders, ages, rank_inc, rank_exc):
    """아이템별 상세 페이지 생성"""
    # ── 제목 박스 ──
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0.1), Inches(10.5), Inches(0.5))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = C_PPT_BLACK
    title_box.line.fill.background()
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    rank_text = ""
    if rank_inc is not None:
        rank_text += f" (전체 {rank_inc}위"
        if rank_exc is not None:
            rank_text += f" / 0제외 {rank_exc}위"
        rank_text += ")"
    r = p.add_run()
    r.text = f"아이템 {item_no}"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = C_PPT_WHITE
    if rank_text:
        r2 = p.add_run()
        r2.text = rank_text
        r2.font.size = Pt(13)
        r2.font.bold = True
        r2.font.color.rgb = RGBColor(0xFF, 0xCC, 0x00)

    # ── 왼쪽: 이미지 영역 ──
    img_left = Inches(0.15)
    img_top = Inches(0.72)
    img_w = Inches(2.2)
    img_h = Inches(2.2)
    img_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_w, img_h)
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = C_LIGHT_GRAY
    img_box.line.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
    img_box.line.width = Pt(1)
    tf_img = img_box.text_frame
    tf_img.word_wrap = True
    tf_img.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_img = tf_img.paragraphs[0]
    p_img.text = "상품 이미지"
    p_img.font.size = Pt(12)
    p_img.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p_img.alignment = PP_ALIGN.CENTER

    # ── 왼쪽 하단: 코멘트 영역 ──
    comment_top = img_top + img_h + Inches(0.08)
    comment_w = img_w
    comment_h = Inches(4.3)
    ch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, comment_top, comment_w, Inches(0.25))
    ch.fill.solid()
    ch.fill.fore_color.rgb = C_NAVY
    ch.line.fill.background()
    tf_ch = ch.text_frame
    tf_ch.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_ch.margin_left = Inches(0.08)
    p_ch = tf_ch.paragraphs[0]
    p_ch.text = "코멘트 / 특이사항"
    p_ch.font.size = Pt(9)
    p_ch.font.bold = True
    p_ch.font.color.rgb = C_PPT_WHITE
    cbox = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, img_left, comment_top + Inches(0.25),
        comment_w, comment_h - Inches(0.25))
    cbox.fill.solid()
    cbox.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xF8)
    cbox.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    cbox.line.width = Pt(1)
    tf_c = cbox.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.08)
    tf_c.margin_top = Inches(0.05)
    p_c = tf_c.paragraphs[0]
    p_c.text = "정성적 의견을 직접 작성해주세요."
    p_c.font.size = Pt(8)
    p_c.font.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
    p_c.font.italic = True

    # ── 데이터 계산 ──
    pref_qs = [q for q in item["questions"] if q["type"] == "선호도"]
    price_qs = [q for q in item["questions"] if q["type"] == "적정가격"]
    color_qs = [q for q in item["questions"] if q["type"] in ("컬러선호", "컬러선호도")]

    pref_avg = calculate_averages(data, pref_qs[0]["col_idx"], gender_col, age_col) if pref_qs else None
    pref_avg_ex = calculate_averages(data, pref_qs[0]["col_idx"], gender_col, age_col, exclude_zero=True) if pref_qs else None
    price_avg = calculate_averages(data, price_qs[0]["col_idx"], gender_col, age_col) if price_qs else None
    price_avg_ex = calculate_averages(data, price_qs[0]["col_idx"], gender_col, age_col, exclude_zero=True) if price_qs else None
    color_stats, color_total = collect_color_stats(data, color_qs[0]["col_idx"], gender_col, age_col) if color_qs else ({"전체": {}}, 0)

    # ── 중앙: 선호도 → 교차분석 → 적정가격 ──
    data_left = 2.50
    data_top = 0.72
    tbl_w = 4.40

    if pref_avg:
        bottom = _add_item_data_table(
            slide, data_left, data_top, tbl_w,
            "선호도 평균 (10점 만점)", C_PPT_BLACK,
            pref_avg, pref_avg_ex, genders, ages, fmt="score")
        data_top = bottom / 914400 + 0.08

    if pref_avg:
        cross_bottom = _add_cross_table(
            slide, data_left, data_top, tbl_w,
            pref_avg, genders, ages, C_PPT_BLACK)
        data_top = cross_bottom / 914400 + 0.08

    if price_avg:
        bottom = _add_item_data_table(
            slide, data_left, data_top, tbl_w,
            "적정가격 평균 (원)", C_PPT_BLACK,
            price_avg, price_avg_ex, genders, ages, fmt="price")
        data_top = bottom / 914400 + 0.08

    # ── 오른쪽: 컬러 선호 ──
    color_left = 7.10
    color_top_in = 0.72
    color_w = 3.55

    if color_stats["전체"]:
        color_bottom = _add_color_table(
            slide, color_left, color_top_in, color_w, C_PPT_BLACK,
            color_stats, color_total)
        color_top = color_bottom / 914400 + 0.08
    else:
        color_top = color_top_in

    # ── 특수 질문 ──
    special_qs = [q for q in item["questions"] if q["type"] not in COMMON_TYPES]
    sq_data = []
    for sq in special_qs:
        stats = collect_special_stats(data, sq["col_idx"], sq["header"], sq["type"])
        if stats["type"] != "empty":
            sq_data.append((sq, stats))

    if sq_data:
        SLIDE_BOTTOM = 7.2
        cur_slide = slide
        sq_left = color_left
        sq_w = color_w
        sq_top = color_top

        def _sq_height(st):
            h = 0.22
            if st["type"] == "numeric":
                h += 0.22
            elif st["type"] == "choice":
                h += 0.18 * (min(len(st["distribution"]), 5) + 1) + 0.08
            return h

        def _make_overflow_slide():
            nonlocal cur_slide, sq_top, sq_left, sq_w
            cur_slide = prs.slides.add_slide(prs.slide_layouts[6])
            bar2 = cur_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0.1), Inches(10.5), Inches(0.5))
            bar2.fill.solid()
            bar2.fill.fore_color.rgb = C_PPT_BLACK
            bar2.line.fill.background()
            tf2 = bar2.text_frame
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf2.margin_left = Inches(0.15)
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = f"아이템 {item_no} - 특수 질문 (계속)"
            r2.font.size = Pt(20)
            r2.font.bold = True
            r2.font.color.rgb = C_PPT_WHITE
            sq_top = 0.80
            sq_left = 0.20
            sq_w = 10.4
            sh = cur_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(sq_left), Inches(sq_top), Inches(sq_w), Inches(0.25))
            sh.fill.solid()
            sh.fill.fore_color.rgb = C_PPT_BLACK
            sh.line.fill.background()
            tf_sh = sh.text_frame
            tf_sh.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf_sh.margin_left = Inches(0.1)
            p_sh = tf_sh.paragraphs[0]
            p_sh.text = "특수 질문 요약 (계속)"
            p_sh.font.size = Pt(10)
            p_sh.font.bold = True
            p_sh.font.color.rgb = C_PPT_WHITE
            sq_top += 0.30

        sh = cur_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(sq_left), Inches(sq_top), Inches(sq_w), Inches(0.25))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_PPT_BLACK
        sh.line.fill.background()
        tf_sh = sh.text_frame
        tf_sh.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_sh.margin_left = Inches(0.1)
        p_sh = tf_sh.paragraphs[0]
        p_sh.text = "특수 질문 요약"
        p_sh.font.size = Pt(10)
        p_sh.font.bold = True
        p_sh.font.color.rgb = C_PPT_WHITE
        sq_top += 0.30

        for sq, stats in sq_data:
            needed = _sq_height(stats)
            if sq_top + needed > SLIDE_BOTTOM:
                _make_overflow_slide()

            q_title = sq["header"]
            if len(q_title) > 50:
                q_title = q_title[:47] + "..."
            tb = cur_slide.shapes.add_textbox(Inches(sq_left), Inches(sq_top), Inches(sq_w), Inches(0.20))
            tf_sq = tb.text_frame
            tf_sq.word_wrap = True
            tf_sq.margin_left = Inches(0.05)
            p_sq = tf_sq.paragraphs[0]
            p_sq.text = f"▸ {q_title}"
            p_sq.font.size = Pt(7.5)
            p_sq.font.bold = True
            p_sq.font.color.rgb = C_NAVY
            sq_top += 0.22

            if stats["type"] == "numeric":
                tb2 = cur_slide.shapes.add_textbox(Inches(sq_left + 0.1), Inches(sq_top), Inches(sq_w - 0.1), Inches(0.18))
                tf2 = tb2.text_frame
                p2 = tf2.paragraphs[0]
                run = p2.add_run()
                run.text = f"평균: {stats['avg']}점 ({stats['count']}명) | 0제외: {stats['avg_ex']}점 ({stats['count_ex']}명)"
                run.font.size = Pt(7)
                run.font.color.rgb = C_PPT_BLACK
                sq_top += 0.22

            elif stats["type"] == "choice":
                top_items = stats["distribution"][:5]
                n_sr = len(top_items) + 1
                tbl_s = cur_slide.shapes.add_table(
                    n_sr, 3, Inches(sq_left), Inches(sq_top),
                    Inches(sq_w), Inches(0.18 * n_sr))
                t = tbl_s.table
                t.first_row = False
                _clear_table_style(tbl_s)
                for i in range(n_sr):
                    t.rows[i].height = Inches(0.18)
                t.columns[0].width = Inches(sq_w * 0.50)
                t.columns[1].width = Inches(sq_w * 0.25)
                t.columns[2].width = Inches(sq_w * 0.25)
                for ci, h in enumerate(["항목", "응답수", "비율"]):
                    c = t.cell(0, ci)
                    c.text = h
                    _style_ppt_cell(c, Pt(6.5), True, C_PPT_WHITE, C_HDR_MAIN)
                    _set_cell_border(c, "000000", "6350", dash=None)
                max_pct = max((round(cnt / stats["total"] * 100, 1) if stats["total"] > 0 else 0) for _, cnt in top_items) if top_items else 0
                for sri, (val, cnt) in enumerate(top_items, 1):
                    pct = round(cnt / stats["total"] * 100, 1) if stats["total"] > 0 else 0
                    is_top = (pct == max_pct and pct > 0)
                    c0 = t.cell(sri, 0)
                    c0.text = val[:30] if len(val) > 30 else val
                    _style_ppt_cell(c0, Pt(6.5), False, C_PPT_BLACK, C_PPT_WHITE, align=PP_ALIGN.LEFT)
                    c0.text_frame.margin_left = Pt(3)
                    _set_cell_border(c0, "000000", "6350")
                    c1 = t.cell(sri, 1)
                    c1.text = str(cnt)
                    _style_ppt_cell(c1, Pt(6.5), is_top, C_SCORE_RED if is_top else C_PPT_BLACK, C_PPT_WHITE)
                    _set_cell_border(c1, "000000", "6350")
                    c2 = t.cell(sri, 2)
                    c2.text = f"{pct}%"
                    _style_ppt_cell(c2, Pt(6.5), is_top, C_SCORE_RED if is_top else C_PPT_BLACK, C_PPT_WHITE)
                    _set_cell_border(c2, "000000", "6350")
                sq_top += 0.18 * n_sr + 0.06


def create_ppt(items, data, gender_col, age_col, genders, ages, output_path):
    """아이템별 1페이지씩 PPT 자동 생성 (새 디자인)"""
    prs = Presentation()
    prs.slide_width = Cm(27.517)
    prs.slide_height = Cm(19.05)

    # ══ 1. 요약 페이지: 선호도 랭킹 (0포함 vs 0제외) ══
    items_inc = []
    items_exc = []
    for item in items:
        pref_qs = [q for q in item["questions"] if q["type"] == "선호도"]
        if not pref_qs:
            continue
        avg = calculate_averages(data, pref_qs[0]["col_idx"], gender_col, age_col)
        avg_ex = calculate_averages(data, pref_qs[0]["col_idx"], gender_col, age_col, exclude_zero=True)
        for target_list, avg_data in [(items_inc, avg), (items_exc, avg_ex)]:
            entry = {
                "item_no": item["item_no"],
                "name": f"아이템 {item['item_no']}",
                "ttl": avg_data["전체"],
                "n_total": avg_data["전체_응답수"],
                "ages": {},
                "n_ages": {},
            }
            for a in ages:
                ad = avg_data["연령"].get(a, {})
                entry["ages"][a] = ad.get("평균", "-")
                entry["n_ages"][a] = ad.get("응답수", 0)
            target_list.append(entry)

    items_inc.sort(key=lambda x: -(x["ttl"] if isinstance(x["ttl"], (int, float)) else 0))
    items_exc.sort(key=lambda x: -(x["ttl"] if isinstance(x["ttl"], (int, float)) else 0))

    all_scores = []
    for it in items_inc + items_exc:
        if isinstance(it["ttl"], (int, float)):
            all_scores.append(it["ttl"])
        for a in ages:
            v = it["ages"].get(a, "-")
            if isinstance(v, (int, float)):
                all_scores.append(v)
    score_min = min(all_scores) if all_scores else 0
    score_max = max(all_scores) if all_scores else 10

    slide_s = prs.slides.add_slide(prs.slide_layouts[6])

    # 타이틀 박스
    title_box = slide_s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0.1), Inches(10.5), Inches(0.5))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = C_PPT_WHITE
    title_box.line.color.rgb = C_PPT_BLACK
    title_box.line.width = Pt(2.5)
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "[유형 선호도] 선호도 랭킹 분석"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_PPT_BLACK

    # 인사이트
    n_total = items_inc[0]["n_total"] if items_inc else 0
    top_names_inc = " > ".join([it["name"] for it in items_inc[:7]])

    tb1 = slide_s.shapes.add_textbox(Inches(0.2), Inches(0.68), Inches(10.4), Inches(0.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    r1a = p1.add_run()
    r1a.text = f"① 전체 모수(n={n_total}): "
    r1a.font.size = Pt(9.5)
    r1a.font.bold = True
    r1a.font.color.rgb = C_PPT_BLACK
    p1b = tf1.add_paragraph()
    p1b.space_before = Pt(1)
    r1c = p1b.add_run()
    r1c.text = f"→ {top_names_inc} 순"
    r1c.font.size = Pt(8.5)
    r1c.font.bold = True
    r1c.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    ranking_changes = []
    for exc_rank, exc_item in enumerate(items_exc):
        inc_rank = next(
            (i for i, inc_item in enumerate(items_inc) if inc_item["item_no"] == exc_item["item_no"]), -1)
        if inc_rank >= 0:
            change = inc_rank - exc_rank
            ranking_changes.append((exc_item["name"], change, exc_rank))

    big_jumpers = [name for name, change, _ in sorted(ranking_changes, key=lambda x: -x[1]) if change >= 3]

    tb2 = slide_s.shapes.add_textbox(Inches(0.2), Inches(1.18), Inches(10.4), Inches(0.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2a = p2.add_run()
    r2a.text = "② 비구매 고객(0점) 제외: "
    r2a.font.size = Pt(9.5)
    r2a.font.bold = True
    r2a.font.color.rgb = C_PPT_BLACK
    r2b = p2.add_run()
    r2b.text = "전체 모수 대비 랭킹이 확 올라간 아이템 → 확실한 선호있는 유형"
    r2b.font.size = Pt(9.5)
    r2b.font.color.rgb = C_PPT_BLACK

    if big_jumpers:
        p2b = tf2.add_paragraph()
        p2b.space_before = Pt(1)
        r2c = p2b.add_run()
        r2c.text = "→ 대폭 상승: "
        r2c.font.size = Pt(8.5)
        r2c.font.bold = True
        r2c.font.color.rgb = C_PPT_BLACK
        for ji, jname in enumerate(big_jumpers[:5]):
            rj = p2b.add_run()
            rj.text = jname
            rj.font.size = Pt(8.5)
            rj.font.bold = True
            rj.font.color.rgb = C_SCORE_RED
            rj.font.underline = True
            if ji < len(big_jumpers[:5]) - 1:
                rsep = p2b.add_run()
                rsep.text = ", "
                rsep.font.size = Pt(8.5)
                rsep.font.color.rgb = C_PPT_BLACK

    # 좌측: 전체 모수 랭킹 테이블
    table_top = 1.78
    tw = 5.20
    _add_ranking_table(slide_s, 0.12, table_top, tw, 5.20,
                       "전체 모수", C_NAVY, items_inc, ages, score_min, score_max)

    # 우측: 비구매 제외 랭킹 테이블
    _add_ranking_table(slide_s, 5.52, table_top, tw, 5.20,
                       "비구매 제외", C_DARK_RED, items_exc, ages, score_min, score_max)

    # ══ 2. 아이템별 상세 페이지 ══
    rank_inc_map = {it["item_no"]: i+1 for i, it in enumerate(items_inc)}
    rank_exc_map = {it["item_no"]: i+1 for i, it in enumerate(items_exc)}

    for item in items:
        item_no = item["item_no"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_item_slide(
            prs, slide, item, item_no, data, gender_col, age_col,
            genders, ages,
            rank_inc_map.get(item_no), rank_exc_map.get(item_no))

    prs.save(output_path)
    print(f"\n📑 PPT 생성 완료! 파일: {output_path}")
    print(f"   - 요약 1페이지 + 아이템 {len(items)}페이지")


# ── PPT 색상 팔레트 ──
C_NAVY = RGBColor(0x2D, 0x3A, 0x4A)
C_DARK_RED = RGBColor(0xA0, 0x00, 0x00)
C_PPT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_PPT_BLACK = RGBColor(0x33, 0x33, 0x33)
C_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
C_BAR_BLUE = RGBColor(0xDD, 0xE8, 0xF6)
C_BAR_GOLD = RGBColor(0xFC, 0xF0, 0xCF)
C_SCORE_RED = RGBColor(0xC0, 0x00, 0x00)
C_HDR_MAIN = RGBColor(0x5B, 0x9B, 0xD5)
C_HDR_EXC = RGBColor(0xC9, 0xB0, 0x7A)
C_TOTAL_BG = RGBColor(0xD6, 0xE4, 0xF0)
C_TOTAL_LABEL = RGBColor(0xB0, 0xCC, 0xE0)


def _style_ppt_cell(cell, font_size, bold, font_color, fill_color, align=PP_ALIGN.CENTER):
    """PPT 테이블 셀 스타일"""
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color
    tf = cell.text_frame
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = False
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        for old in bodyPr.findall(qn('a:normAutofit')) + bodyPr.findall(qn('a:spAutoFit')) + bodyPr.findall(qn('a:noAutofit')):
            bodyPr.remove(old)
        etree.SubElement(bodyPr, qn('a:normAutofit'), fontScale='100000', lnSpcReduction='0')
    for p in tf.paragraphs:
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _set_cell_border(cell, color_hex="000000", width="6350", dash="dot"):
    """셀 테두리 설정 (기본: 검정 0.5pt 점선)"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge_tag in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = tcPr.find(qn(edge_tag))
        if ln is not None:
            tcPr.remove(ln)
        ln = etree.SubElement(tcPr, qn(edge_tag), w=width, cmpd="sng")
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr"), val=color_hex)
        if dash:
            etree.SubElement(ln, qn("a:prstDash"), val=dash)


def _clear_table_style(tbl_shape):
    """테이블 기본 테마 스타일 제거 (테두리 덮어쓰기 방지)"""
    tbl_xml = tbl_shape._element
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = tbl_xml.find('.//' + qn('a:tblPr'))
    if tblPr is not None:
        for attr in ['bandRow', 'bandCol', 'firstRow', 'lastRow', 'firstCol', 'lastCol']:
            if attr in tblPr.attrib:
                del tblPr.attrib[attr]
        for ts in tblPr.findall(qn('a:tblStyle')):
            tblPr.remove(ts)
        for child in list(tblPr):
            if 'tblStyle' in child.tag:
                tblPr.remove(child)


def _data_bar_width_ratio(val, min_val, max_val):
    """값 → 0.0~1.0 비율 (데이터 막대 너비용)"""
    if not isinstance(val, (int, float)):
        return 0.0
    if max_val <= min_val:
        return 0.5
    return max(0.0, min(1.0, (val - min_val) / (max_val - min_val)))


def _apply_cell_data_bar(cell, ratio, bar_color, bg_color=RGBColor(0xFF, 0xFF, 0xFF)):
    """셀 내부에 그라디언트 채우기로 데이터 막대 효과 적용"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old_fill in tcPr.findall(qn("a:solidFill")) + tcPr.findall(qn("a:gradFill")) + tcPr.findall(qn("a:noFill")):
        tcPr.remove(old_fill)
    if ratio <= 0.01:
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        etree.SubElement(solidFill, qn("a:srgbClr"), val=f"{bg_color}")
        return
    ratio = min(ratio, 1.0)
    bar_end = int(ratio * 100000)
    bar_end_next = min(bar_end + 1, 100000)
    bar_hex = f"{bar_color[0]:02X}{bar_color[1]:02X}{bar_color[2]:02X}"
    bg_hex = f"{bg_color[0]:02X}{bg_color[1]:02X}{bg_color[2]:02X}"
    gradFill = etree.SubElement(tcPr, qn("a:gradFill"), rotWithShape="0")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    gs1 = etree.SubElement(gsLst, qn("a:gs"), pos="0")
    etree.SubElement(gs1, qn("a:srgbClr"), val=bar_hex)
    gs2 = etree.SubElement(gsLst, qn("a:gs"), pos=str(bar_end))
    etree.SubElement(gs2, qn("a:srgbClr"), val=bar_hex)
    gs3 = etree.SubElement(gsLst, qn("a:gs"), pos=str(bar_end_next))
    etree.SubElement(gs3, qn("a:srgbClr"), val=bg_hex)
    gs4 = etree.SubElement(gsLst, qn("a:gs"), pos="100000")
    etree.SubElement(gs4, qn("a:srgbClr"), val=bg_hex)
    etree.SubElement(gradFill, qn("a:lin"), ang="0", scaled="0")


def main():
    if len(sys.argv) < 2:
        # 기본 파일 경로
        input_file = "RAW 파일.xlsx"
    else:
        input_file = sys.argv[1]

    if not os.path.exists(input_file):
        print(f"❌ 파일을 찾을 수 없습니다: {input_file}")
        sys.exit(1)

    # 출력 파일명 생성
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    output_excel = f"{base_name}_분석결과_v7.xlsx"
    output_ppt = f"{base_name}_분석결과_v7.pptx"

    print(f"📊 설문조사 분석 시작: {input_file}")
    print("=" * 50)

    # 1. 데이터 로드
    print("1️⃣  데이터 로딩 중...")
    headers, data = load_raw_data(input_file)
    print(f"   → {len(data)}명 응답자, {len(headers)}개 열")

    # 2. 성별/연령 열 찾기
    gender_col, age_col = find_gender_age_columns(headers)
    print(f"2️⃣  성별 열: {headers[gender_col][:30] if gender_col else '없음'}")
    print(f"   연령 열: {headers[age_col][:30] if age_col else '없음'}")

    # 3. 아이템 식별
    print("3️⃣  아이템 식별 중...")
    items = identify_items(headers)
    print(f"   → {len(items)}개 아이템 발견")
    for item in items:
        q_types = [q["type"] for q in item["questions"]]
        print(f"      아이템 {item['item_no']}: {len(item['questions'])}개 질문 ({', '.join(set(q_types))})")

    # 성별/연령 목록
    genders_set = set()
    ages_set = set()
    for row in data:
        g = clean_gender(row[gender_col]) if gender_col is not None and row[gender_col] else "미응답"
        a = str(row[age_col]).strip() if age_col is not None and row[age_col] else "미응답"
        genders_set.add(g)
        ages_set.add(a)
    genders = sorted(genders_set)
    ages = sort_age_groups(ages_set)

    # 4. 결과 엑셀 생성
    print("4️⃣  분석 결과 엑셀 생성 중...")
    create_summary_excel(headers, data, items, gender_col, age_col, output_excel)

    # 5. PPT 생성
    print("5️⃣  PPT 생성 중...")
    create_ppt(items, data, gender_col, age_col, genders, ages, output_ppt)

    print("\n" + "=" * 50)
    print(f"🎉 모든 작업 완료!")
    print(f"   📗 엑셀: {output_excel}")
    print(f"   📙 PPT:  {output_ppt}")


def run_analysis_from_bytes(file_bytes, filename="upload.xlsx"):
    """웹 업로드용: 바이트 데이터를 받아 분석 후 (excel_bytes, ppt_bytes) 반환"""
    import io, tempfile, shutil

    tmpdir = tempfile.mkdtemp()
    try:
        input_path = os.path.join(tmpdir, filename)
        with open(input_path, "wb") as f:
            f.write(file_bytes)

        base_name = os.path.splitext(filename)[0]
        excel_path = os.path.join(tmpdir, f"{base_name}_분석결과.xlsx")
        ppt_path = os.path.join(tmpdir, f"{base_name}_분석결과.pptx")

        headers, data = load_raw_data(input_path)
        gender_col, age_col = find_gender_age_columns(headers)
        items = identify_items(headers)

        genders_set = set()
        ages_set = set()
        for row in data:
            g = clean_gender(row[gender_col]) if gender_col is not None and row[gender_col] else "미응답"
            a = str(row[age_col]).strip() if age_col is not None and row[age_col] else "미응답"
            genders_set.add(g)
            ages_set.add(a)
        genders = sorted(genders_set)
        ages = sort_age_groups(ages_set)

        create_summary_excel(headers, data, items, gender_col, age_col, excel_path)
        create_ppt(items, data, gender_col, age_col, genders, ages, ppt_path)

        with open(excel_path, "rb") as f:
            excel_bytes = f.read()
        with open(ppt_path, "rb") as f:
            ppt_bytes = f.read()

        return {
            "excel_bytes": excel_bytes,
            "ppt_bytes": ppt_bytes,
            "excel_name": f"{base_name}_분석결과.xlsx",
            "ppt_name": f"{base_name}_분석결과.pptx",
            "n_respondents": len(data),
            "n_items": len(items),
            "n_headers": len(headers),
            "items": items,
        }
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


if __name__ == "__main__":
    main()
